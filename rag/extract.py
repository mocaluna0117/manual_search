"""PDF以外のファイルから本文を取り出す。

対応形式ごとに「区切りの単位」が違うので、どれも (見出し, 本文) の並びに
揃えて返す。この見出しがそのまま引用の手がかり(PDFのページ番号に相当)になる。

- Word(.docx)   … 段落と表。区切りが無いので全体で1つ
- Excel(.xlsx)  … シートごと
- PowerPoint(.pptx) … スライドごと
- Outlook(.msg) … 差出人・宛先・件名・本文をまとめて1つ

いずれも「読めない部分は飛ばして、読めた分だけ返す」方針。1か所の不備で
ファイル全体を取り込めなくするより、部分的にでも検索できる方が役に立つ。
"""

from io import BytesIO

# 1セルが極端に長い場合の打ち切り(壊れたファイルでメモリを食い潰さないため)
MAX_CELL_CHARS = 2000


def extract_docx(data: bytes) -> list[tuple[str, str]]:
    """Word文書から本文を取り出す。段落と表の中身を上から順に拾う"""
    import docx

    document = docx.Document(BytesIO(data))
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]

    # 表は行ごとに「セル | セル」の形にする(読み下せる並びを保つ)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    return [("", text)] if text else []


def extract_xlsx(data: bytes) -> list[tuple[str, str]]:
    """Excelブックから本文を取り出す。シート単位で分け、行を1行ずつ並べる。

    read_only + data_only にするのが要点。data_only=Trueで数式ではなく
    計算結果の値を読む(「=SUM(...)」ではなく実際の数字を検索対象にする)
    """
    import openpyxl

    book = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
    sections: list[tuple[str, str]] = []
    try:
        for sheet in book.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [
                    str(v).strip()[:MAX_CELL_CHARS]
                    for v in row
                    if v is not None and str(v).strip()
                ]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sections.append((sheet.title, "\n".join(rows)))
    finally:
        book.close()
    return sections


def extract_pptx(data: bytes) -> list[tuple[str, str]]:
    """PowerPointから本文を取り出す。スライド単位で分け、発表者ノートも拾う"""
    from pptx import Presentation

    presentation = Presentation(BytesIO(data))
    sections: list[tuple[str, str]] = []
    for i, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            # 表は行ごとに読み下す
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        # 発表者ノートには手順の補足が書かれていることが多い
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"(ノート) {notes}")
        if parts:
            sections.append((f"スライド{i}", "\n".join(parts)))
    return sections


def extract_msg(data: bytes) -> list[tuple[str, str]]:
    """Outlookのメール(.msg)から差出人・宛先・件名・本文を取り出す。

    問い合わせの回答メールは、本文だけでなく件名や日付も手がかりになるので
    ヘッダも本文に含めて検索できるようにする。
    """
    import extract_msg as msg_lib

    message = msg_lib.Message(BytesIO(data))
    try:
        header_lines = [
            f"件名: {message.subject}" if message.subject else "",
            f"差出人: {message.sender}" if message.sender else "",
            f"宛先: {message.to}" if message.to else "",
            f"日付: {message.date}" if message.date else "",
        ]
        body = (message.body or "").strip()
        attachments = [
            a.longFilename or a.shortFilename
            for a in (message.attachments or [])
            if (a.longFilename or a.shortFilename)
        ]
        if attachments:
            header_lines.append(f"添付: {', '.join(attachments)}")
        text = "\n".join([h for h in header_lines if h] + ["", body]).strip()
    finally:
        message.close()
    return [("", text)] if text else []


# 拡張子から抽出処理を引く表。PDFは元からある経路で扱うのでここには載せない
EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".msg": extract_msg,
}


def extractor_for(file_name: str):
    """ファイル名に対応する抽出処理を返す(未対応ならNone)"""
    lower = (file_name or "").lower()
    for ext, fn in EXTRACTORS.items():
        if lower.endswith(ext):
            return fn
    return None
